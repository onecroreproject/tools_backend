from fastapi import FastAPI, UploadFile, File, Form, APIRouter
from pydantic import BaseModel, EmailStr
import bcrypt
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
import shutil, os, uuid, zipfile, datetime
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
import pikepdf
from pdf2image import convert_from_path
import shutil
import uuid
import re
from pathlib import Path
import pdfkit
import time
from PIL import Image, ImageDraw, ImageFont
from database import files_collection
from database import users_collection
from database import otp_collection
import ocrmypdf
from pptx import Presentation
import fitz  
from pdf2docx import Converter
import pdfplumber
import pandas as pd
from docx import Document
# import comtypes.client
# from docx2pdf import convert
import subprocess
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.pagesizes import A4
from fastapi import HTTPException
import secrets
import string
import random 
import io
from rembg import remove
from email_utils import send_otp_email
import yt_dlp
import tempfile
from typing import Optional
import instaloader
import requests
from bs4 import BeautifulSoup
import subprocess
from pathlib import Path
import imgkit
import speedtest
# from moviepy.editor import VideoFileClip
import platform
import segno
from fastapi.middleware.cors import CORSMiddleware


app = FastAPI(title="PDF Tools")

origins = [
    "http://localhost:3000"
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,  # or ["*"] to allow all (not recommended in prod)
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

TEMP_DIR = "temp"
ENDPOINTS = ["merge","split","compress","pdf-to-image","crop", "rotate", "png-to-jpg", "jpg-to-png", "watermark"]
for endpoint in ENDPOINTS:
    os.makedirs(os.path.join(TEMP_DIR,endpoint,"uploads"),exist_ok=True)
    os.makedirs(os.path.join(TEMP_DIR,endpoint,"output"),exist_ok=True)

POPPLER_PATH = r"C:\Program Files\Release-25.07.0-0\poppler-25.07.0\Library\bin"
soffice_path = r"C:\Program Files\LibreOffice\program\soffice.exe"

# Helper → Save metadata to MongoDB
async def save_metadata(operation: str, filename: str, size: int):
    doc = {
        "operation": operation,
        "filename": filename,
        "size": size,
        "timestamp": datetime.datetime.utcnow()
    }
    await files_collection.insert_one(doc)

def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(plain: str, hashed: str) -> bool:
    return bcrypt.checkpw(plain.encode('utf-8'), hashed.encode('utf-8'))

# Pydantic models
class SignupModel(BaseModel):
    name: str
    email: EmailStr
    phone_number: str
    password: str

class LoginModel(BaseModel):
    email: EmailStr
    password: str

class ForgotPasswordModel(BaseModel):
    email: EmailStr
    new_password: str
    confirm_password: str

class OTPRequestModel(BaseModel):
    email: EmailStr

class OTPVerifyModel(BaseModel):
    otp: str

class ResetPasswordModel(BaseModel):
    new_password: str
    confirm_password: str

def generate_otp(length=6) -> str:
    return ''.join(random.choices(string.digits, k=length))

async def save_otp(email: str, otp: str):
    await otp_collection.delete_many({"email": email})  # remove previous OTPs
    await otp_collection.insert_one({
        "email": email,
        "otp": otp,
        "expires_at": datetime.datetime.utcnow() + datetime.timedelta(minutes=10)
    })

async def verify_otp(email: str, otp: str) -> bool:
    record = await otp_collection.find_one({"email": email})
    if record and record["otp"] == otp:
        if record["expires_at"] > datetime.datetime.utcnow():
            return True  # ✅ VALID OTP
    return False  # ❌ INVALID or EXPIRED


@app.post("/signup/request-otp")
async def request_signup_otp(user: SignupModel):
    try:
        existing_user = await users_collection.find_one({"email": user.email})
        if existing_user:
            raise Exception("Email already registered.")

        otp = generate_otp()

        result = await otp_collection.update_one(
            {"email": user.email},
            {"$set": {
                "otp": otp,
                "expires_at": datetime.datetime.utcnow() + datetime.timedelta(minutes=10),
                "name": user.name,
                "password": hash_password(user.password)
            }},
            upsert=True
        )

        print(f"OTP Update Result: {result.raw_result}")

        saved = await otp_collection.find_one({"email": user.email})
        print(f"Document after saving OTP: {saved}")

        await send_otp_email(user.email, otp, "Signup")
        return {"message": "OTP sent to your email for verification."}

    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/signup/verify-otp")
async def verify_signup_otp(data: OTPVerifyModel):
    try:
        record = await otp_collection.find_one({"otp": data.otp})
        if not record:
            raise Exception("Invalid or expired OTP.")

        if record["expires_at"] < datetime.datetime.utcnow():
            raise Exception("OTP expired.")

        await users_collection.insert_one({
            "name": record["name"],
            "email": record["email"],
            "hashed_password": record["password"]
        })

        await otp_collection.delete_one({"otp": data.otp})
        return {"message": "User registered successfully."}

    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

# Login endpoint
@app.post("/login")
async def login(user: LoginModel):
    try:
        existing_user = await users_collection.find_one({"email": user.email})
        if not existing_user:
            raise Exception("User not found.")

        if not verify_password(user.password, existing_user["hashed_password"]):
            raise Exception("Invalid password.")

        return {"message": "Login successful."}

    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.post("/forgot-password/request-otp")
async def forgot_password_request(data: OTPRequestModel):
    try:
        user = await users_collection.find_one({"email": data.email})
        if not user:
            raise Exception("User not found.")

        otp = generate_otp()

        await otp_collection.update_one(
            {"email": data.email},
            {
                "$set": {
                    "otp": otp,
                    "expires_at": datetime.datetime.utcnow() + datetime.timedelta(minutes=10),
                    "otp_verified": False
                }
            },
            upsert=True
        )

        await send_otp_email(data.email, otp, "Password Reset")
        return {"message": "OTP sent to your email."}

    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/forgot-password/verify-otp")
async def verify_forgot_password(data: OTPVerifyModel):
    try:
        record = await otp_collection.find_one({"otp": data.otp})
        if not record:
            raise Exception("Invalid OTP.")

        if record["expires_at"] < datetime.datetime.utcnow():
            raise Exception("OTP expired.")

        await otp_collection.update_one(
            {"otp": data.otp},
            {"$set": {"otp_verified": True}}
        )

        return {"message": "OTP verified. You can now reset your password."}

    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))

@app.post("/forgot-password/reset")
async def reset_password(data: ResetPasswordModel):
    try:
        if data.new_password != data.confirm_password:
            raise Exception("Passwords do not match.")

        # Find the verified record
        record = await otp_collection.find_one({"otp_verified": True})
        if not record:
            raise Exception("No verified OTP session found. Please verify OTP again.")

        await users_collection.update_one(
            {"email": record["email"]},
            {"$set": {"hashed_password": hash_password(data.new_password)}}
        )

        # Clean up the OTP entry
        await otp_collection.delete_one({"email": record["email"]})

        return {"message": "Password reset successful."}

    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


# ✅ 1. Merge PDFs
@app.post("/merge")
async def merge_pdfs(files: list[UploadFile] = File(...)):

    upload_dir = os.path.join(TEMP_DIR,"merge","uploads")
    output_dir = os.path.join(TEMP_DIR,"merge","output")
    merger = PdfMerger()

    for file in files:
        temp_path = os.path.join(upload_dir,f"{uuid.uuid4().hex}_{file.filename}")
        with open(temp_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        merger.append(temp_path)

    output_file = os.path.join(output_dir,f"merged_{uuid.uuid4().hex}.pdf")
    merger.write(output_file)
    merger.close()

    await save_metadata("merge",os.path.basename(output_file), os.path.getsize(output_file))
    return FileResponse(output_file, media_type="application/pdf", filename="merged.pdf")


# ✅ 2. Split PDF
@app.post("/split")
async def split_pdf(file: UploadFile = File(...), pages: str = Form(...)):
    import re

    upload_dir = os.path.join(TEMP_DIR, "split", "uploads")
    output_dir = os.path.join(TEMP_DIR, "split", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    reader = PdfReader(input_path)
    total_pages = len(reader.pages)

    # Clean and validate input
    pages = pages.replace(" ", "")
    if not re.fullmatch(r"([0-9]+-[0-9]+)(,[0-9]+-[0-9]+)*", pages):
        return {"error": "Invalid format. Use ranges like 1-4,5-8"}

    zip_filename = os.path.join(output_dir, f"split_{uuid.uuid4().hex}.zip")

    with zipfile.ZipFile(zip_filename, "w") as zipf:
        for idx, r in enumerate(pages.split(","), start=1):
            start, end = map(int, r.split("-"))

            if start < 1 or end > total_pages or start > end:
                return {"error": f"Invalid range: {start}-{end}. Total pages = {total_pages}"}

            writer = PdfWriter()
            for i in range(start - 1, end):
                writer.add_page(reader.pages[i])

            split_filename = f"part_{idx}_{start}-{end}.pdf"
            split_path = os.path.join(output_dir, split_filename)

            with open(split_path, "wb") as f:
                writer.write(f)

            zipf.write(split_path, arcname=split_filename)
            os.remove(split_path)

    await save_metadata("split", os.path.basename(zip_filename), os.path.getsize(zip_filename))
    return FileResponse(zip_filename, media_type="application/zip", filename="split_parts.zip")


# ✅ 3. Compress PDF
@app.post("/compress")
async def compress_pdf(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR,"compress","uploads")
    output_dir = os.path.join(TEMP_DIR,"compress","output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir,f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    output_path = os.path.join(output_dir, f"compress_{uuid.uuid4().hex}.pdf")

    try:
        gs_command = [
            "gswin64c",
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.4",
            "-dPDFSETTINGS=/ebook",  # Change to /screen or /printer as needed
            "-dNOPAUSE",
            "-dQUIET",
            "-dBATCH",
            f"-sOutputFile={output_path}",
            input_path
        ]
        subprocess.run(gs_command, check=True)
    except subprocess.CalledProcessError as e:
        return {"error": f"Ghostscript compression failed: {str(e)}"}

    await save_metadata("compress", os.path.basename(output_path), os.path.getsize(output_path))
    return FileResponse(output_path, media_type="application/pdf", filename="compressed.pdf")

# ✅ 4. PDF to Images
@app.post("/pdf-to-image")
async def pdf_to_image(file: UploadFile = File(...)):

    upload_dir = os.path.join(TEMP_DIR,"pdf-to-image","uploads")
    output_dir = os.path.join(TEMP_DIR,"pdf-to-image","output")


    input_path = os.path.join(upload_dir,f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        # Convert PDF to images
        images = convert_from_path(input_path, dpi=200, poppler_path=POPPLER_PATH)
    except Exception as e:
        return {"error": f"PDF to image conversion failed: {str(e)}"}

    # Create a zip file to store images
    zip_path = os.path.join(output_dir, f"images_{uuid.uuid4().hex}.zip")

    with zipfile.ZipFile(zip_path, "w") as zipf:
        for i, img in enumerate(images, start=1):
            img_path = os.path.join(output_dir, f"page_{i}.jpg")
            img.save(img_path, "JPEG")
            zipf.write(img_path, f"page_{i}.jpg")
            os.remove(img_path)

    # Optionally log or save metadata
    await save_metadata("pdf-to-image", os.path.basename(zip_path), os.path.getsize(zip_path))
    # Return zip file
    return FileResponse(zip_path, media_type="application/zip", filename="pdf_images.zip")
    
@app.post("/remove-pages")
async def remove_pages(file: UploadFile = File(...), pages: str = Form(...)):
    import re

    upload_dir = os.path.join(TEMP_DIR, "remove-pages", "uploads")
    output_dir = os.path.join(TEMP_DIR, "remove-pages", "output")

    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    reader = PdfReader(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)

    pages = pages.replace(" ", "")
    if not re.fullmatch(r"[0-9,-]+", pages):
        return {"error": "Invalid page format. Use numbers or ranges like 1,2-5"}

    remove_indices = set()
    for r in pages.split(","):
        if "-" in r:
            start, end = map(int, r.split("-"))
            remove_indices.update(range(start - 1, end))
        else:
            remove_indices.add(int(r) - 1)

    if any(i >= total_pages or i < 0 for i in remove_indices):
        return {"error": "Page numbers out of range."}

    for i in range(total_pages):
        if i not in remove_indices:
            writer.add_page(reader.pages[i])

    output_file = os.path.join(output_dir, f"removed_{uuid.uuid4().hex}.pdf")
    with open(output_file, "wb") as f:
        writer.write(f)

    await save_metadata("remove-pages", os.path.basename(output_file), os.path.getsize(output_file))
    return FileResponse(output_file, media_type="application/pdf", filename="pages_removed.pdf")

@app.post("/extract-pages")
async def extract_pages(file: UploadFile = File(...), pages: str = Form(...)):

    upload_dir = os.path.join(TEMP_DIR, "extract-pages", "uploads")
    output_dir = os.path.join(TEMP_DIR, "extract-pages", "output")

    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    reader = PdfReader(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)

    pages = pages.replace(" ", "")
    if not re.fullmatch(r"[0-9,-]+", pages):
        return {"error": "Invalid page format. Use numbers or ranges like 1,3-5"}

    extract_indices = []
    for r in pages.split(","):
        if "-" in r:
            start, end = map(int, r.split("-"))
            extract_indices.extend(range(start - 1, end))
        else:
            extract_indices.append(int(r) - 1)

    if any(i >= total_pages or i < 0 for i in extract_indices):
        return {"error": "Page numbers out of range."}

    for i in extract_indices:
        writer.add_page(reader.pages[i])

    output_file = os.path.join(output_dir, f"extracted_{uuid.uuid4().hex}.pdf")
    with open(output_file, "wb") as f:
        writer.write(f)

    await save_metadata("extract-pages", os.path.basename(output_file), os.path.getsize(output_file))
    return FileResponse(output_file, media_type="application/pdf", filename="extracted_pages.pdf")

@app.post("/organize")
async def organize_pdf(file: UploadFile = File(...), order: str = Form(...)):
    import re

    upload_dir = os.path.join(TEMP_DIR, "organize", "uploads")
    output_dir = os.path.join(TEMP_DIR, "organize", "output")

    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    reader = PdfReader(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)

    order = order.replace(" ", "")
    if not re.fullmatch(r"[0-9,]+", order):
        return {"error": "Invalid order format. Use comma-separated numbers like 3,1,2"}

    try:
        page_indices = [int(i) - 1 for i in order.split(",")]
    except ValueError:
        return {"error": "Order must be numeric."}

    if len(page_indices) != total_pages:
        return {"error": "Order must include all pages exactly once."}
    if sorted(page_indices) != list(range(total_pages)):
        return {"error": "Invalid page indices. Must use all pages without duplicates."}

    for i in page_indices:
        writer.add_page(reader.pages[i])

    output_file = os.path.join(output_dir, f"organized_{uuid.uuid4().hex}.pdf")
    with open(output_file, "wb") as f:
        writer.write(f)

    await save_metadata("organize", os.path.basename(output_file), os.path.getsize(output_file))
    return FileResponse(output_file, media_type="application/pdf", filename="organized.pdf")

@app.post("/repair")
async def repair_pdf(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "repair", "uploads")
    output_dir = os.path.join(TEMP_DIR, "repair", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    output_file = os.path.join(output_dir, f"repaired_{uuid.uuid4().hex}.pdf")

    # --- Try PikePDF first ---
    try:
        pdf = pikepdf.open(input_path, allow_overwriting_input=True)
        pdf.save(output_file, fix_metadata=True, linearize=True)
        pdf.close()
    except Exception as e:
        # --- Fall back to PyPDF2 ---
        try:
            from PyPDF2 import PdfReader, PdfWriter
            reader = PdfReader(input_path, strict=False)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            with open(output_file, "wb") as f:
                writer.write(f)
        except Exception:
            return {"error": f"Failed to repair PDF: {str(e)}"}

    await save_metadata("repair", os.path.basename(output_file), os.path.getsize(output_file))
    return FileResponse(output_file, media_type="application/pdf", filename="repaired.pdf")


@app.post("/ocr")
async def ocr_pdf(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "ocr", "uploads")
    output_dir = os.path.join(TEMP_DIR, "ocr", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    output_file = os.path.join(output_dir, f"ocr_{uuid.uuid4().hex}.pdf")
    try:
        ocrmypdf.ocr(input_path, output_file, deskew=True, optimize=3, force_ocr=True)
    except Exception as e:
        return {"error": f"OCR failed: {str(e)}"}

    await save_metadata("ocr", os.path.basename(output_file), os.path.getsize(output_file))
    return FileResponse(output_file, media_type="application/pdf", filename="ocr.pdf")


@app.post("/jpg-to-pdf")
async def jpg_to_pdf(files: list[UploadFile] = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "jpg-to-pdf", "uploads")
    output_dir = os.path.join(TEMP_DIR, "jpg-to-pdf", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    image_list = []
    for file in files:
        input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
        with open(input_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        img = Image.open(input_path).convert("RGB")
        image_list.append(img)

    output_file = os.path.join(output_dir, f"jpg2pdf_{uuid.uuid4().hex}.pdf")
    image_list[0].save(output_file, save_all=True, append_images=image_list[1:])

    await save_metadata("jpg-to-pdf", os.path.basename(output_file), os.path.getsize(output_file))
    return FileResponse(output_file, media_type="application/pdf", filename="images.pdf")


@app.post("/html-to-pdf")
async def html_to_pdf(file: UploadFile = File(...)):
    # Define directories
    upload_dir = os.path.join(TEMP_DIR, "html-to-pdf", "uploads")
    output_dir = os.path.join(TEMP_DIR, "html-to-pdf", "output")

    # Ensure the directories exist
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    # Save the uploaded HTML file
    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # Define the output PDF file path
    output_file = os.path.join(output_dir, f"html2pdf_{uuid.uuid4().hex}.pdf")
    
    try:
        # Convert HTML file to PDF using pdfkit
        pdfkit.from_file(input_path, output_file)
    except Exception as e:
        return {"error": f"Failed to convert HTML: {str(e)}"}

    # Save metadata to the database (optional)
    await save_metadata("html-to-pdf", os.path.basename(output_file), os.path.getsize(output_file))

    # Return the generated PDF file
    return FileResponse(output_file, media_type="application/pdf", filename="html.pdf")


@app.post("/pdf-to-ppt")
async def pdf_to_ppt(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "pdf-to-ppt", "uploads")
    output_dir = os.path.join(TEMP_DIR, "pdf-to-ppt", "output")

    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # Convert PDF pages to images
    try:
        # adjust dpi for better quality, maybe 200–300
        images = convert_from_path(input_path, dpi=200)
    except Exception as e:
        return {"error": f"Failed converting PDF to images: {str(e)}"}

    # Create PowerPoint presentation
    prs = Presentation()

    blank_slide_layout = prs.slide_layouts[6]  # blank layout

    for i, pil_img in enumerate(images):
        slide = prs.slides.add_slide(blank_slide_layout)

        # Save the page image temporarily
        img_filename = os.path.join(upload_dir, f"page_{i}.png")
        pil_img.save(img_filename, format="PNG")

        # Get image dimensions in pixels
        width_px, height_px = pil_img.size

        # Convert to inches: assume 96 dpi if needed, or use other
        # But pptx uses EMUs: use helper to scale
        # Simplest: fill the entire slide with the image

        slide.shapes.add_picture(img_filename, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # Optionally remove the temp image file
        os.remove(img_filename)

    output_file = os.path.join(output_dir, f"pdf_to_pptx_{uuid.uuid4().hex}.pptx")
    prs.save(output_file)

    return FileResponse(output_file, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                          filename="presentation.pptx")

# @app.post("/pdf-to-word")
# async def pdf_to_word(file: UploadFile = File(...)):
#     upload_dir = os.path.join(TEMP_DIR, "pdf-to-word", "uploads")
#     output_dir = os.path.join(TEMP_DIR, "pdf-to-word", "output")
    
#     os.makedirs(upload_dir, exist_ok=True)
#     os.makedirs(output_dir, exist_ok=True)

#     input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
#     with open(input_path, "wb") as buffer:
#         shutil.copyfileobj(file.file, buffer)

#     output_file = os.path.join(output_dir, f"pdf_to_word_{uuid.uuid4().hex}.docx")

#     try:
#         # Convert PDF to Word using pdf2docx
#         converter = Converter(input_path)
#         converter.convert(output_file, start=0, end=None)  # Convert entire document
#         converter.close()
        
#         await save_metadata("pdf-to-word", os.path.basename(output_file), os.path.getsize(output_file))
#         return FileResponse(output_file, media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", filename="document.docx")
#     except Exception as e:
#         return {"error": f"PDF to Word conversion failed: {str(e)}"}


@app.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "pdf-to-excel", "uploads")
    output_dir = os.path.join(TEMP_DIR, "pdf-to-excel", "output")

    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    output_file = os.path.join(output_dir, f"pdf_to_excel_{uuid.uuid4().hex}.xlsx")

    try:
        with pdfplumber.open(input_path) as pdf:
            with pd.ExcelWriter(output_file, engine='openpyxl') as writer:
                sheet_count = 1
                wrote_sheet = False  # Track if at least one sheet was written

                for page in pdf.pages:
                    table = page.extract_table()
                    if table:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        df.to_excel(writer, sheet_name=f"Sheet_{sheet_count}", index=False)
                        sheet_count += 1
                        wrote_sheet = True

            # Ensure at least one sheet was written
            if not wrote_sheet:
                return {"error": "No tables were found in the PDF to write to Excel."}

        return FileResponse(
            output_file,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename="document.xlsx"
        )

    except Exception as e:
        return {"error": f"PDF to Excel conversion failed: {str(e)}"}


@app.post("/word-to-pdf")
async def word_to_pdf(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "word-to-pdf", "uploads")
    output_dir = os.path.join(TEMP_DIR, "word-to-pdf", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        subprocess.run([
            soffice_path,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            input_path
        ], check=True)

        # Construct the output PDF path
        output_filename = os.path.splitext(os.path.basename(input_path))[0] + ".pdf"
        output_file = os.path.join(output_dir, output_filename)

        if not os.path.exists(output_file):
            return {"error": "Conversion failed: Output PDF not found."}

        await save_metadata("word-to-pdf", output_filename, os.path.getsize(output_file))

        return FileResponse(output_file, media_type="application/pdf", filename="document.pdf")

    except subprocess.CalledProcessError as e:
        return {"error": f"LibreOffice conversion failed: {e}"}
    except Exception as e:
        return {"error": f"Unexpected error: {str(e)}"}

@app.post("/ppt-to-pdf")
async def ppt_to_pdf(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "ppt-to-pdf", "uploads")
    output_dir = os.path.join(TEMP_DIR, "ppt-to-pdf", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        # Use LibreOffice in headless mode
        subprocess.run([
            soffice_path,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", output_dir,
            input_path
        ], check=True)

        # The output file will have the same base name as input but with .pdf extension
        base_filename = os.path.splitext(os.path.basename(input_path))[0]
        output_file = os.path.join(output_dir, f"{base_filename}.pdf")

        if not os.path.exists(output_file):
            return {"error": "Conversion failed: Output PDF not found."}

        await save_metadata("ppt-to-pdf", os.path.basename(output_file), os.path.getsize(output_file))

        return FileResponse(output_file, media_type="application/pdf", filename="presentation.pdf")

    except subprocess.CalledProcessError as e:
        return {"error": f"LibreOffice conversion failed: {e}"}
    except Exception as e:
        return {"error": f"Unexpected error: {str(e)}"}


@app.post("/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "excel-to-pdf", "uploads")
    output_dir = os.path.join(TEMP_DIR, "excel-to-pdf", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = Path(upload_dir) / f"{uuid.uuid4().hex}.xlsx"
    input_path = input_path.resolve()
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        # Run LibreOffice CLI to convert Excel to PDF
        subprocess.run([
            soffice_path,  # or full path: "C:\\Program Files\\LibreOffice\\program\\soffice.exe"
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(output_dir),
            str(input_path)
        ], check=True)

        # Find the output file
        output_file = output_dir + "/" + input_path.stem + ".pdf"
        output_file = Path(output_file)

        if not output_file.exists():
            return {"error": "Conversion failed. PDF not created."}

        await save_metadata("excel-to-pdf", output_file.name, os.path.getsize(output_file))

        return FileResponse(str(output_file), media_type="application/pdf", filename="document.pdf")

    except subprocess.CalledProcessError as e:
        return {"error": f"LibreOffice conversion failed: {str(e)}"}

@app.post("/rotate-pdf")
async def rotate_pdf(file: UploadFile = File(...), angle: int = Form(...)):
    upload_dir = os.path.join(TEMP_DIR, "rotate-pdf", "uploads")
    output_dir = os.path.join(TEMP_DIR, "rotate-pdf", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    output_file = os.path.join(output_dir, f"rotated_{uuid.uuid4().hex}.pdf")

    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()

        # Rotate each page
        for page in reader.pages:
            page.rotate(angle)
            writer.add_page(page)

        with open(output_file, "wb") as f:
            writer.write(f)

    except Exception as e:
        return {"error": f"Failed to rotate PDF: {str(e)}"}

    await save_metadata("rotate-pdf", os.path.basename(output_file), os.path.getsize(output_file))

    return FileResponse(output_file, media_type="application/pdf", filename="rotated.pdf")

@app.post("/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "pdf-to-excel", "uploads")
    output_dir = os.path.join(TEMP_DIR, "pdf-to-excel", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    output_file = os.path.join(output_dir, f"pdf_to_excel_{uuid.uuid4().hex}.xlsx")

    try:
        # Extract tables from PDF using Tabula
        tables = tabula.read_pdf(input_path, pages='all', multiple_tables=True)

        # Create an Excel writer to save multiple DataFrames
        with pd.ExcelWriter(output_file) as writer:
            for i, table in enumerate(tables):
                table.to_excel(writer, sheet_name=f"Sheet_{i + 1}", index=False)

    except Exception as e:
        return {"error": f"Failed to convert PDF to Excel: {str(e)}"}

    await save_metadata("pdf-to-excel", os.path.basename(output_file), os.path.getsize(output_file))

    return FileResponse(output_file, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", filename="document.xlsx")

@app.post("/add-watermark")
async def add_watermark(file: UploadFile = File(...), watermark_text: str = Form(...)):

    upload_dir = os.path.join(TEMP_DIR, "add-watermark", "uploads")
    output_dir = os.path.join(TEMP_DIR, "add-watermark", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    watermark_pdf_path = os.path.join(upload_dir, "watermark.pdf")
    output_path = os.path.join(output_dir, f"watermarked_{uuid.uuid4().hex}.pdf")

    # Save uploaded file
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # Create watermark PDF
    c = canvas.Canvas(watermark_pdf_path, pagesize=letter)
    c.setFont("Helvetica", 40)
    c.setFillAlpha(0.3)
    c.drawCentredString(300, 500, watermark_text)
    c.save()

    # Merge watermark with each page
    reader = PdfReader(input_path)
    writer = PdfWriter()
    watermark = PdfReader(watermark_pdf_path).pages[0]

    for page in reader.pages:
        page.merge_page(watermark)
        writer.add_page(page)

    with open(output_path, "wb") as f:
        writer.write(f)

    await save_metadata("add-watermark", os.path.basename(output_path), os.path.getsize(output_path))
    return FileResponse(output_path, media_type="application/pdf", filename="watermarked.pdf")


@app.post("/remove-watermark")
async def remove_watermark(
    file: UploadFile = File(...),
    keywords: str = Form("dinesh,watermark,draft,confidential")
):
    """
    Removes text, image, and vector watermarks from a PDF safely.
    """

    upload_dir = os.path.join(TEMP_DIR, "remove-watermark", "uploads")
    output_dir = os.path.join(TEMP_DIR, "remove-watermark", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    output_path = os.path.join(output_dir, f"cleaned_{uuid.uuid4().hex}.pdf")

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        doc = fitz.open(input_path)
        words_to_remove = [k.strip().lower() for k in keywords.split(",") if k.strip()]

        for page in doc:
            # 1️⃣ Remove text watermarks
            for word in words_to_remove:
                areas = page.search_for(word)
                for area in areas:
                    page.add_redact_annot(area, fill=(1, 1, 1))
            page.apply_redactions()


            drawings = page.get_drawings()
            for d in drawings:
                try:
                    # Transparent filled rectangles
                    if "fill" in d and d["fill"] and len(d["fill"]) == 4 and d["fill"][3] < 0.5:
                        page.add_redact_annot(d["rect"], fill=(1, 1, 1))

                    # Check line items
                    if "items" in d:
                        for item in d["items"]:
                            if isinstance(item, (list, tuple)) and len(item) >= 2 and item[0] == "l":
                                coords = item[1]
                                if len(coords) == 4:
                                    x0, y0, x1, y1 = coords
                                    # Diagonal line → likely watermark
                                    if abs(x1 - x0) > 200 and abs(y1 - y0) > 200:
                                        page.add_redact_annot(d["rect"], fill=(1, 1, 1))
                except Exception:
                    continue
            page.apply_redactions()

        doc.save(output_path)
        doc.close()

    except Exception as e:
        return {"error": f"Watermark removal failed: {str(e)}"}

    await save_metadata("remove-watermark", os.path.basename(output_path), os.path.getsize(output_path))
    return FileResponse(output_path, media_type="application/pdf", filename="no_watermark.pdf")


@app.post("/add-page-numbers")
async def add_page_numbers(file: UploadFile = File(...)):
    from PyPDF2 import PdfReader, PdfWriter

    upload_dir = os.path.join(TEMP_DIR, "add-page-numbers", "uploads")
    output_dir = os.path.join(TEMP_DIR, "add-page-numbers", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    reader = PdfReader(input_path)
    writer = PdfWriter()
    total_pages = len(reader.pages)

    for i, page in enumerate(reader.pages, start=1):
        # Create overlay with page number
        packet_path = os.path.join(upload_dir, f"number_{i}.pdf")
        c = canvas.Canvas(packet_path, pagesize=letter)
        c.setFont("Helvetica", 12)
        c.drawString(500, 10, f"{i} / {total_pages}")
        c.save()

        overlay = PdfReader(packet_path).pages[0]
        page.merge_page(overlay)
        writer.add_page(page)
        os.remove(packet_path)

    output_path = os.path.join(output_dir, f"paged_{uuid.uuid4().hex}.pdf")
    with open(output_path, "wb") as f:
        writer.write(f)

    await save_metadata("add-page-numbers", os.path.basename(output_path), os.path.getsize(output_path))
    return FileResponse(output_path, media_type="application/pdf", filename="page_numbered.pdf")

@app.post("/crop-pdf")
async def crop_pdf(file: UploadFile = File(...), crop_box: str = Form(...)):
    """
    crop_box format: x0,y0,x1,y1 (in points)
    Example: "50,50,500,700"
    """
    upload_dir = os.path.join(TEMP_DIR, "crop-pdf", "uploads")
    output_dir = os.path.join(TEMP_DIR, "crop-pdf", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    output_path = os.path.join(output_dir, f"cropped_{uuid.uuid4().hex}.pdf")

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        # Validate crop_box format
        try:
            x0, y0, x1, y1 = map(float, crop_box.split(","))
        except ValueError:
            return {"error": "Invalid crop_box format. Expected 'x0,y0,x1,y1' with numeric values, e.g., '50,50,500,700'"}

        reader = PdfReader(input_path)
        writer = PdfWriter()

        for page in reader.pages:
            page.cropbox.lower_left = (x0, y0)
            page.cropbox.upper_right = (x1, y1)
            writer.add_page(page)

        with open(output_path, "wb") as f:
            writer.write(f)

    except Exception as e:
        return {"error": f"Failed to crop PDF: {str(e)}"}

    await save_metadata("crop-pdf", os.path.basename(output_path), os.path.getsize(output_path))
    return FileResponse(output_path, media_type="application/pdf", filename="cropped.pdf")

@app.post("/edit-pdf-text")
async def edit_pdf(file: UploadFile = File(...), old_text: str = Form(...), new_text: str = Form(...)):
    upload_dir = os.path.join(TEMP_DIR, "edit-pdf", "uploads")
    output_dir = os.path.join(TEMP_DIR, "edit-pdf", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    output_path = os.path.join(output_dir, f"edited_{uuid.uuid4().hex}.pdf")

    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        doc = fitz.open(input_path)
        for page in doc:
            areas = page.search_for(old_text)
            for area in areas:
                page.add_redact_annot(area, fill=(1, 1, 1))
            page.apply_redactions()
            for area in areas:
                page.insert_text(area.tl, new_text, fontsize=12)

        doc.save(output_path)
    except Exception as e:
        return {"error": f"Edit PDF failed: {str(e)}"}

    await save_metadata("edit-pdf", os.path.basename(output_path), os.path.getsize(output_path))
    return FileResponse(output_path, media_type="application/pdf", filename="edited.pdf")

@app.post("/unlock")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form(...)):
    """
    Remove a known user/owner password from a PDF.
    Provide the correct password via form.
    """
    upload_dir = os.path.join(TEMP_DIR, "unlock", "uploads")
    output_dir = os.path.join(TEMP_DIR, "unlock", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    output_file = os.path.join(output_dir, f"unlocked_{uuid.uuid4().hex}.pdf")

    # Try pikepdf first (fast & reliable)
    try:
        pdf = pikepdf.open(input_path, password=password)
        # save without password (unlocked)
        pdf.save(output_file)
        pdf.close()
    except Exception as pike_err:
        # Fallback to PyPDF2 if pikepdf fails
        try:
            reader = PdfReader(input_path)
            # For older PyPDF2 versions, .decrypt may be required:
            try:
                # PyPDF2 >= 2.x
                if reader.is_encrypted:
                    reader.decrypt(password)
            except Exception:
                # attempt decrypt via PdfReader constructor if supported
                reader = PdfReader(input_path, password=password)

            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            with open(output_file, "wb") as out_f:
                writer.write(out_f)
        except Exception as py_err:
            return {"error": f"Failed to unlock PDF. pikepdf error: {pike_err} ; PyPDF2 error: {py_err}"}

    await save_metadata("unlock", os.path.basename(output_file), os.path.getsize(output_file))
    return FileResponse(output_file, media_type="application/pdf", filename="unlocked.pdf")


@app.post("/protect")
async def protect_pdf(
    file: UploadFile = File(...),
    user_password: str = Form(...),
    owner_password: str = Form(None),
    encryption_level: int = Form(4)  # R value, 4 is standard (AES-128); you can expose if needed
):

    upload_dir = os.path.join(TEMP_DIR, "protect", "uploads")
    output_dir = os.path.join(TEMP_DIR, "protect", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    output_file = os.path.join(output_dir, f"protected_{uuid.uuid4().hex}.pdf")

    try:
        owner_pwd = owner_password if owner_password else user_password
        pdf = pikepdf.open(input_path)
        enc = pikepdf.Encryption(user=user_password, owner=owner_pwd, R=encryption_level)
        pdf.save(output_file, encryption=enc)
        pdf.close()
    except Exception as e:
        # Try PyPDF2 fallback (simpler encryption)
        try:
            reader = PdfReader(input_path)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            # PyPDF2.encrypt signature may vary: try common variants
            try:
                writer.encrypt(user_password, owner_pwd)
            except TypeError:
                # Some PyPDF2 versions accept only one arg
                writer.encrypt(user_password)
            with open(output_file, "wb") as out_f:
                writer.write(out_f)
        except Exception as py_err:
            return {"error": f"Failed to protect PDF: {e} ; PyPDF2 fallback error: {py_err}"}

    await save_metadata("protect", os.path.basename(output_file), os.path.getsize(output_file))
    return FileResponse(output_file, media_type="application/pdf", filename="protected.pdf")

@app.post("/resize-image")
async def resize_image(file: UploadFile = File(...), width: int = Form(...), height: int = Form(...)):
    upload_dir = os.path.join(TEMP_DIR, "resize", "uploads")
    output_dir = os.path.join(TEMP_DIR, "resize", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        # Open the image
        img = Image.open(input_path)
        
        # Check if the image has an alpha channel (RGBA)
        if img.mode == 'RGBA':
            # Create a new image with a white background and paste the RGBA image onto it
            img = img.convert("RGB")  # Converts to RGB, removing alpha channel

        # Resize the image
        img = img.resize((width, height))

        # Save the resized image
        resized_path = os.path.join(output_dir, f"resized_{uuid.uuid4().hex}.jpg")
        img.save(resized_path, "JPEG")

    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Image resize failed: {str(e)}")

    return FileResponse(resized_path, media_type="image/jpeg", filename="resized_image.jpg")

@app.post("/scale-image")
async def scale_image(file: UploadFile = File(...), factor: float = Form(...)):
    upload_dir = os.path.join(TEMP_DIR, "scale", "uploads")
    output_dir = os.path.join(TEMP_DIR, "scale", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        # Open and convert image
        img = Image.open(input_path)
        if img.mode in ("RGBA", "LA", "P"):
            img = img.convert("RGB")

        # Calculate new size
        width, height = img.size
        new_width = int(width * factor)
        new_height = int(height * factor)

        # Use high-quality resampling filter
        img = img.resize((new_width, new_height), resample=Image.LANCZOS)

        # Save output
        scaled_path = os.path.join(output_dir, f"scaled_{uuid.uuid4().hex}.jpg")
        img.save(scaled_path, "JPEG", quality=95, optimize=True)

    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Image scaling failed: {str(e)}")

    return FileResponse(scaled_path, media_type="image/jpeg", filename="scaled_image.jpg")

#password generator
def generate_password(length: int, use_uppercase: bool, use_lowercase: bool, use_numbers: bool, use_symbols: bool) -> str:
    """
    Generate a secure password based on the selected complexity options.
    """
    # Base sets of characters
    character_sets = []

    if use_uppercase:
        character_sets.append(string.ascii_uppercase)
    if use_lowercase:
        character_sets.append(string.ascii_lowercase)
    if use_numbers:
        character_sets.append(string.digits)
    if use_symbols:
        character_sets.append(string.punctuation)

    # If no character set is selected, raise an exception
    if not character_sets:
        raise HTTPException(status_code=400, detail="At least one character set must be selected")

    # Generate the password by randomly selecting characters from the chosen sets
    password = ''.join(random.choice(''.join(character_sets)) for _ in range(length))

    return password


# Password generator endpoint
@app.post("/generate-password")
async def generate_password_endpoint(
    length: int = Form(12),  # Default password length is 12
    use_uppercase: bool = Form(True),  # Default includes uppercase
    use_lowercase: bool = Form(True),  # Default includes lowercase
    use_numbers: bool = Form(True),  # Default includes numbers
    use_symbols: bool = Form(True)  # Default includes symbols
):
    """
    Endpoint to generate a password based on specified criteria.
    """
    # Validate length (must be between 8 and 32 characters)
    if length < 8 or length > 64:
        raise HTTPException(status_code=400, detail="Password length must be between 8 and 32 characters.")

    # Generate password using the provided parameters
    password = generate_password(length, use_uppercase, use_lowercase, use_numbers, use_symbols)

    # Return the generated password as a JSON response
    return JSONResponse(content={"password": password})

@app.post("/compress-image")
async def compress_image(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "compress-image", "uploads")
    output_dir = os.path.join(TEMP_DIR, "compress-image", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    # Output path for compressed image
    output_path_jpeg = os.path.join(output_dir, f"compressed_{uuid.uuid4().hex}.jpg")
    output_path_png = os.path.join(output_dir, f"compressed_{uuid.uuid4().hex}.png")

    try:
        # Open the image
        img = Image.open(input_path)

        # If the image has an alpha channel (RGBA), we need to convert it to RGB for JPEG
        if img.mode == "RGBA":
            img = img.convert("RGB")
            img.save(output_path_jpeg, format="JPEG", quality=30, optimize=True)
            return FileResponse(output_path_jpeg, media_type="image/jpeg", filename="compressed_image.jpg")
        
        # If it's in RGB mode or any other format, we can directly compress as JPEG
        if img.mode in ["RGB", "L"]:
            img.save(output_path_jpeg, format="JPEG", quality=30, optimize=True)
            return FileResponse(output_path_jpeg, media_type="image/jpeg", filename="compressed_image.jpg")

        # If transparency is needed (RGBA), save as PNG
        img.save(output_path_png, format="PNG", optimize=True)
        return FileResponse(output_path_png, media_type="image/png", filename="compressed_image.png")
        
    except Exception as e:
        return {"error": f"Image compression failed: {str(e)}"}

@app.post("/remove-background")
async def remove_background(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "remove-background", "uploads")
    output_dir = os.path.join(TEMP_DIR, "remove-background", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    output_path = os.path.join(output_dir, f"no_bg_{uuid.uuid4().hex}.png")

    try:
        # Open the image file and remove background
        with open(input_path, "rb") as input_file:
            input_data = input_file.read()
            output_data = remove(input_data)

        # Save the output data as a PNG file (since the output will have transparent background)
        with open(output_path, "wb") as output_file:
            output_file.write(output_data)

    except Exception as e:
        return {"error": f"Background removal failed: {str(e)}"}

    await save_metadata("remove-background", os.path.basename(output_path), os.path.getsize(output_path))
    return FileResponse(output_path, media_type="image/png", filename="image_no_bg.png")

# Directory setup for YouTube and Instagram downloads
DOWNLOAD_DIRS = {
    "youtube": os.path.join(TEMP_DIR, "youtube", "output"),
    "instagram": os.path.join(TEMP_DIR, "instagram", "output"),
}

# Make the output folders if they don't exist
for path in DOWNLOAD_DIRS.values():
    os.makedirs(path, exist_ok=True)

@app.post("/download/youtube")
async def download_youtube(
    url: str = Form(...),
    mode: str = Form("video"),  # "video" or "audio"
    resolution: Optional[str] = Form(None)  # e.g., "720", "1080", "2160", "4320"
):
    try:
        output_dir = DOWNLOAD_DIRS["youtube"]
        outtmpl = os.path.join(output_dir, '%(title)s.%(ext)s')

        # Validate and parse resolution
        res_int = None
        if resolution:
            try:
                res_int = int(resolution)
            except ValueError:
                raise HTTPException(status_code=400, detail="Invalid resolution format. Use numeric string like '720'.")

        # Construct format string based on mode and resolution
        if mode == "audio":
            ydl_format = "bestaudio"
        elif res_int:
            # video + audio with video height <= resolution
            ydl_format = f"bestvideo[height<={res_int}]+bestaudio/best[height<={res_int}]"
        else:
            ydl_format = "bestvideo+bestaudio/best"

        ydl_opts = {
            'format': ydl_format,
            'outtmpl': outtmpl,
            'merge_output_format': 'mp4',
            'quiet': True,
            'cookiefile' : r"R:\dinesh\tools\pdf_tools\youtube.com_cookies.txt"
        }

        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            filename = ydl.prepare_filename(info)

        # Fix filename extension for audio/video formats
        if mode == "audio":
            # yt-dlp might save audio in formats other than mp3
            # If needed, convert or rename extension to mp3
            base, ext = os.path.splitext(filename)
            if ext not in [".mp3", ".m4a", ".webm", ".opus"]:
                filename = base + ".mp3"
        else:
            # Replace webm extension with mp4 for video if applicable
            filename = filename.replace(".webm", ".mp4")

        if not os.path.exists(filename):
            raise HTTPException(status_code=404, detail="Download failed or file not found.")

        await save_metadata(
            f"youtube_{mode}_download",
            os.path.basename(filename),
            os.path.getsize(filename)
        )

        media_type = "audio/mpeg" if mode == "audio" else "video/mp4"

        return FileResponse(
            filename,
            filename=os.path.basename(filename),
            media_type=media_type
        )

    except yt_dlp.utils.DownloadError as e:
        # Provide clearer message if YouTube blocks download due to login required
        if "Sign in to confirm" in str(e):
            raise HTTPException(
                status_code=403,
                detail="YouTube video requires login and cannot be downloaded without authentication."
            )
        else:
            raise HTTPException(status_code=400, detail=f"YouTube {mode} download failed: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"YouTube {mode} download failed: {str(e)}")


@app.post("/download/instagram")
async def download_instagram(url: str = Form(...)):
    output_dir = DOWNLOAD_DIRS["instagram"]
    outtmpl = os.path.join(output_dir, '%(title)s.%(ext)s')

    # First, try downloading video or reel using yt-dlp
    try:
        ydl_opts = {
            'format': 'best',
            'outtmpl': outtmpl,
            'quiet': True,
            'cookiefile': r"R:\dinesh\tools\pdf_tools\youtube.com_cookies",  # ✅ Adjust this path if needed
            'noplaylist': True,
        }

        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            info = ydl.extract_info(url, download=True)
            filename = ydl.prepare_filename(info)

            if not os.path.exists(filename):
                raise Exception("File not found after yt-dlp download.")

            # Return video/reel file
            return FileResponse(
                path=filename,
                filename=os.path.basename(filename),
                media_type="video/mp4"
            )

    except Exception as video_error:
        error_message = str(video_error)

        # Fallback to instaloader only if it's a non-video post
        if "There is no video in this post" in error_message or "Unsupported URL" in error_message:
            try:
                shortcode = url.strip("/").split("/")[-1]

                # Temporary unique folder for the post
                post_dir = os.path.join(output_dir, str(uuid.uuid4()))
                os.makedirs(post_dir, exist_ok=True)

                loader = instaloader.Instaloader(
                    dirname_pattern=post_dir,
                    save_metadata=False,
                    download_comments=False,
                    quiet=True
                )

                # 🔐 OPTIONAL: Load session for private post access
                # loader.load_session_from_file("your_username", "session-instagram")

                post = instaloader.Post.from_shortcode(loader.context, shortcode)
                loader.download_post(post, target="")

                # Look for downloaded image
                for file in os.listdir(post_dir):
                    if file.lower().endswith(('.jpg', '.jpeg', '.png', '.webp')):
                        image_path = os.path.join(post_dir, file)
                        return FileResponse(
                            path=image_path,
                            filename=file,
                            media_type="image/jpeg"
                        )

                raise HTTPException(status_code=404, detail="Image not found after download.")

            except Exception as image_error:
                raise HTTPException(status_code=400, detail=f"Image download failed: {str(image_error)}")

        # If it's another yt-dlp error not related to media type, return it
        raise HTTPException(status_code=400, detail=f"Instagram download failed: {error_message}")

@app.post("/image/crop")
async def crop_image(
    file: UploadFile = File(...),
    left: int = Form(...),
    top: int = Form(...),
    right: int = Form(...),
    bottom: int = Form(...)
):
    upload_dir = os.path.join(TEMP_DIR, "crop", "uploads")
    output_dir = os.path.join(TEMP_DIR, "crop", "output")

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    image = Image.open(input_path)
    cropped = image.crop((left, top, right, bottom))

    output_path = os.path.join(output_dir, f"cropped_{uuid.uuid4().hex}.png")
    cropped.save(output_path)

    return FileResponse(output_path, media_type="image/png", filename="cropped.png")

@app.post("/image/rotate")
async def rotate_image(file: UploadFile = File(...), angle: int = Form(...)):
    upload_dir = os.path.join(TEMP_DIR, "rotate", "uploads")
    output_dir = os.path.join(TEMP_DIR, "rotate", "output")

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    image = Image.open(input_path)
    rotated = image.rotate(angle, expand=True)

    output_path = os.path.join(output_dir, f"rotated_{uuid.uuid4().hex}.png")
    rotated.save(output_path)

    return FileResponse(output_path, media_type="image/png", filename="rotated.png")

@app.post("/image/png-to-jpg")
async def png_to_jpg(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "png-to-jpg", "uploads")
    output_dir = os.path.join(TEMP_DIR, "png-to-jpg", "output")

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    image = Image.open(input_path).convert("RGB")

    output_path = os.path.join(output_dir, f"converted_{uuid.uuid4().hex}.jpg")
    image.save(output_path, format="JPEG")

    return FileResponse(output_path, media_type="image/jpeg", filename="converted.jpg")

@app.post("/image/jpg-to-png")
async def jpg_to_png(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "jpg-to-png", "uploads")
    output_dir = os.path.join(TEMP_DIR, "jpg-to-png", "output")

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    image = Image.open(input_path)

    output_path = os.path.join(output_dir, f"converted_{uuid.uuid4().hex}.png")
    image.save(output_path, format="PNG")

    return FileResponse(output_path, media_type="image/png", filename="converted.png")

@app.post("/image/watermark")
async def watermark_image(
    file: UploadFile = File(...),
    watermark_text: str = Form(...),
    position: str = Form(default="bottom-right")
):
    upload_dir = os.path.join(TEMP_DIR, "watermark", "uploads")
    output_dir = os.path.join(TEMP_DIR, "watermark", "output")

    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as f:
        shutil.copyfileobj(file.file, f)

    image = Image.open(input_path).convert("RGBA")

    txt = Image.new("RGBA", image.size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(txt)

    font_size = int(min(image.size) * 0.05)
    font_path = str(Path("arial.ttf"))  # Use your actual path if needed
    font = ImageFont.truetype(font_path, font_size)

    bbox = draw.textbbox((0, 0), watermark_text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]


    positions = {
        "top-left": (10, 10),
        "bottom-right": (image.width - text_width - 10, image.height - text_height - 10),
        "center": ((image.width - text_width) // 2, (image.height - text_height) // 2),
    }

    pos = positions.get(position.lower(), positions["bottom-right"])

    draw.text(pos, watermark_text, font=font, fill=(0, 0, 0, 180))

    watermarked = Image.alpha_composite(image, txt)

    output_path = os.path.join(output_dir, f"watermarked_{uuid.uuid4().hex}.png")
    watermarked.convert("RGB").save(output_path, "PNG")

    return FileResponse(output_path, media_type="image/png", filename="watermarked.png")

@app.post("/html-file-to-image")
async def html_file_to_image(file: UploadFile = File(...)):
    upload_dir = os.path.join(TEMP_DIR, "html-to-image", "uploads")
    output_dir = os.path.join(TEMP_DIR, "html-to-image", "output")

    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)

    # Save uploaded HTML file
    input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
    with open(input_path, "wb") as f:
        f.write(await file.read())

    # Prepare output image path
    output_path = os.path.join(output_dir, f"{uuid.uuid4().hex}.png")

    try:
        imgkit.from_file(input_path, output_path)
    except Exception as e:
        return {"error": f"Conversion failed: {str(e)}"}

    return FileResponse(output_path, media_type="image/png", filename="converted.png")


def draw_text(draw, text, position, font, image_width):
    text = text.upper()
    bbox = draw.textbbox((0, 0), text, font=font)
    text_width = bbox[2] - bbox[0]
    text_height = bbox[3] - bbox[1]
    x = (image_width - text_width) / 2
    y = position

    outline_range = 2
    for ox in range(-outline_range, outline_range + 1):
        for oy in range(-outline_range, outline_range + 1):
            draw.text((x + ox, y + oy), text, font=font, fill="black")

    draw.text((x, y), text, font=font, fill="white")


@app.post("/generate_meme")
async def generate_meme(
    image: UploadFile = File(...),
    top_text: str = Form(""),
    bottom_text: str = Form("")
):
    UPLOAD_DIR = os.path.join(TEMP_DIR, "meme", "uploads")
    OUTPUT_DIR = os.path.join(TEMP_DIR, "meme", "output")

    os.makedirs(UPLOAD_DIR, exist_ok=True)
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    # Save the uploaded image
    input_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4().hex}_{image.filename}")
    with open(input_path, "wb") as buffer:
        buffer.write(await image.read())

    # Open image with PIL
    img = Image.open(input_path)
    draw = ImageDraw.Draw(img)
    image_width, image_height = img.size

    # Choose font and size (adjust path to font if needed)
    font_path = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"
    font_size = int(image_height / 10)
    font = ImageFont.truetype(font_path, font_size)

    # Draw top text
    if top_text:
        draw_text(draw, top_text, 10, font, image_width)

    # Draw bottom text
    if bottom_text:
        bottom_y = image_height - font_size - 10
        draw_text(draw, bottom_text, bottom_y, font, image_width)

    # Save output
    output_filename = f"meme_{uuid.uuid4().hex}.png"
    output_path = os.path.join(OUTPUT_DIR, output_filename)
    img.save(output_path)

    # You can implement save_metadata like in your example if needed
    # await save_metadata("meme", output_filename, os.path.getsize(output_path))

    return FileResponse(output_path, media_type="image/png", filename=output_filename)

@app.get("/speedtest")
def run_speedtest():
    st = speedtest.Speedtest()
    st.get_best_server()
    download_speed = st.download()
    upload_speed = st.upload()
    ping = st.results.ping

    return {
        "download_mbps": round(download_speed / 1_000_000, 2),
        "upload_mbps": round(upload_speed / 1_000_000, 2),
        "ping_ms": round(ping, 2)
    }

# @app.post("/convert_video_format")
# async def convert_video_format(
#     file: UploadFile = File(...),
#     target_format: str = "mp4"
# ):
#     upload_dir = os.path.join(TEMP_DIR, "video_format_convert", "uploads")
#     output_dir = os.path.join(TEMP_DIR, "video_format_convert", "output")
#     os.makedirs(upload_dir, exist_ok=True)
#     os.makedirs(output_dir, exist_ok=True)

#     input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
#     with open(input_path, "wb") as buffer:
#         buffer.write(await file.read())

#     clip = VideoFileClip(input_path)
#     output_filename = f"converted_{uuid.uuid4().hex}.{target_format}"
#     output_path = os.path.join(output_dir, output_filename)

#     clip.write_videofile(output_path, codec="libx264")
#     clip.close()

#     return FileResponse(output_path, media_type="video/mp4", filename=output_filename)

# @app.post("/video_to_gif")
# async def video_to_gif(file: UploadFile = File(...)):

#     upload_dir = os.path.join(TEMP_DIR, "video-to-gif", "uploads")
#     output_dir = os.path.join(TEMP_DIR, "video-to-gif", "output")
#     os.makedirs(upload_dir, exist_ok=True)
#     os.makedirs(output_dir, exist_ok=True)

#     input_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
#     with open(input_path, "wb") as buffer:
#         buffer.write(await file.read())

#     clip = VideoFileClip(input_path)
#     gif_path = os.path.join(output_dir, f"gif_{uuid.uuid4().hex}.gif")
#     clip.write_gif(gif_path)
#     clip.close()

#     return FileResponse(gif_path, media_type="image/gif", filename="converted.gif")

@app.post("/images_to_gif")
async def images_to_gif(files: list[UploadFile] = File(...)):

    upload_dir = os.path.join(TEMP_DIR, "images-to-gif", "uploads")
    output_dir = os.path.join(TEMP_DIR, "images-to-gif", "output")
    os.makedirs(upload_dir, exist_ok=True)
    os.makedirs(output_dir, exist_ok=True)
    
    images = []
    for file in files:
        img_path = os.path.join(upload_dir, f"{uuid.uuid4().hex}_{file.filename}")
        with open(img_path, "wb") as buffer:
            buffer.write(await file.read())
        img = Image.open(img_path).convert("RGBA")
        images.append(img)

    output_path = os.path.join(output_dir, f"animated_{uuid.uuid4().hex}.gif")
    images[0].save(
        output_path,
        save_all=True,
        append_images=images[1:],
        duration=500,
        loop=0
    )

    return FileResponse(output_path, media_type="image/gif", filename="animated.gif")

@app.post("/ping")
async def ping_host(
    host: str = Form(...)
):
    system = platform.system().lower()

    # Use -n for Windows, -c for Unix-based systems
    count_flag = "-n" if system == "windows" else "-c"
    command = ["ping", count_flag, "4", host]

    try:
        result = subprocess.run(command, capture_output=True, text=True, timeout=10)
        output = result.stdout

        if result.returncode != 0:
            return JSONResponse(status_code=400, content={"error": "Ping failed", "details": output})

        # Parse output for latency
        latency_info = parse_ping_output(output, system)
        return latency_info

    except subproc%ss.TimeoutExpired:
        return JSONResponse(status_code=408, content={"error": "Ping request timed out"})


def parse_ping_output(output: str, system: str):
    lines = output.splitlines()
    stats = {}

    if system == "windows":
        for line in lines:
            if "Average" in line:
                # Example: Minimum = 10ms, Maximum = 12ms, Average = 11ms
                parts = line.split(",")
                for part in parts:
                    key, value = part.strip().split(" = ")
                    stats[key.lower()] = value
    else:
        # Unix-based systems
        for line in lines:
            if "min/avg/max" in line or "rtt min/avg/max" in line:
                # Example: rtt min/avg/max/mdev = 10.234/11.567/12.901/0.899 ms
                numbers = line.split("=")[1].strip().split(" ")[0]
                min_, avg, max_, *_ = numbers.split("/")
                stats["min"] = f"{min_} ms"
                stats["avg"] = f"{avg} ms"
                stats["max"] = f"{max_} ms"

    return {
        "host": stats,
        "raw_output": output
    }

@app.post("/qr/generate")
async def generate_qr(
    data: str = Form(...),
    fg_color: str = Form(default="black"),
    bg_color: str = Form(default="white"),
    scale: int = Form(default=10),  # controls resolution
    border: int = Form(default=4),  # controls quiet zone
    logo: UploadFile = File(default=None)
):
    # Output path setup
    output_dir = os.path.join(TEMP_DIR, "qr-generator", "output")
    os.makedirs(output_dir, exist_ok=True)

    file_name = f"qr_{uuid.uuid4().hex}.png"
    file_path = os.path.join(output_dir, file_name)

    # Generate QR code with segno
    qr = segno.make(data, error='h')  # High error correction

    # Save QR code to buffer with given resolution and colors
    qr_buffer = io.BytesIO()
    qr.save(qr_buffer, kind='png', scale=scale, border=border, dark=fg_color, light=bg_color)
    qr_buffer.seek(0)
    qr_img = Image.open(qr_buffer).convert("RGBA")

    # If logo is provided
    if logo is not None:
        logo_bytes = await logo.read()
        logo_img = Image.open(io.BytesIO(logo_bytes)).convert("RGBA")

        # Resize logo proportionally to QR size (e.g., 20%)
        qr_width, qr_height = qr_img.size
        logo_scale_factor = 5  # 1/5th of QR size
        logo_size = (qr_width // logo_scale_factor, qr_height // logo_scale_factor)

        # Resize with high-quality resampling
        logo_img = logo_img.resize(logo_size, Image.Resampling.LANCZOS)

        # Optional: Add rounded corners to logo
        corner_radius = int(min(logo_size) * 0.2)  # 20% radius
        rounded_logo = Image.new("RGBA", logo_size, (0, 0, 0, 0))
        mask = Image.new("L", logo_size, 0)
        draw = ImageDraw.Draw(mask)
        draw.rounded_rectangle([(0, 0), logo_size], corner_radius, fill=255)
        rounded_logo.paste(logo_img, (0, 0), mask=mask)

        # Center the logo
        pos = ((qr_width - logo_size[0]) // 2, (qr_height - logo_size[1]) // 2)
        qr_img.paste(rounded_logo, pos, mask=rounded_logo)

    # Save final image
    qr_img.save(file_path)

    return FileResponse(file_path, media_type="image/png", filename="qr_code.png")


###
length_units = {
    "mm": 0.001,
    "cm": 0.01,
    "m": 1,
    "km": 1000,
    "inch": 0.0254,
    "ft": 0.3048,
    "yd": 0.9144,
    "mile": 1609.34
}

# -----------------------------
# Weight Conversion (Base: gram)
weight_units = {
    "mg": 0.001,
    "g": 1,
    "kg": 1000,
    "ton": 1_000_000,
    "oz": 28.3495,
    "lb": 453.592
}

# -----------------------------
# Time Conversion (Base: seconds)
time_units = {
    "seconds": 1,
    "minutes": 60,
    "hours": 3600,
    "days": 86400
}

# -----------------------------
# Temperature requires formula
def convert_temperature(value, from_unit, to_unit):
    if from_unit == to_unit:
        return value
    if from_unit == "celsius":
        if to_unit == "fahrenheit":
            return (value * 9/5) + 32
        elif to_unit == "kelvin":
            return value + 273.15
    elif from_unit == "fahrenheit":
        if to_unit == "celsius":
            return (value - 32) * 5/9
        elif to_unit == "kelvin":
            return (value - 32) * 5/9 + 273.15
    elif from_unit == "kelvin":
        if to_unit == "celsius":
            return value - 273.15
        elif to_unit == "fahrenheit":
            return (value - 273.15) * 9/5 + 32
    raise ValueError("Invalid temperature conversion")

# -----------------------------
@app.post("/convert_units")
async def convert(
    value: float = Form(...),
    from_unit: str = Form(...),
    to_unit: str = Form(...)
):
    from_unit = from_unit.lower()
    to_unit = to_unit.lower()

    try:
        # Length
        if from_unit in length_units and to_unit in length_units:
            base_value = value * length_units[from_unit]
            converted = base_value / length_units[to_unit]
        # Weight
        elif from_unit in weight_units and to_unit in weight_units:
            base_value = value * weight_units[from_unit]
            converted = base_value / weight_units[to_unit]
        # Time
        elif from_unit in time_units and to_unit in time_units:
            base_value = value * time_units[from_unit]
            converted = base_value / time_units[to_unit]
        # Temperature
        elif from_unit in ["celsius", "fahrenheit", "kelvin"] and to_unit in ["celsius", "fahrenheit", "kelvin"]:
            converted = convert_temperature(value, from_unit, to_unit)
        else:
            return JSONResponse(status_code=400, content={"error": "Unsupported unit conversion"})

        return {
            "input": f"{value} {from_unit}",
            "output": f"{converted:.4f} {to_unit}"
        }

    except Exception as e:
        return JSONResponse(status_code=500, content={"error": str(e)})
    
